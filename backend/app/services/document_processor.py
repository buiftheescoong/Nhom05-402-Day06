import logging
import re
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import settings

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Unstructured-based smart parser (PDF, DOCX, PPTX, TXT …)
# ---------------------------------------------------------------------------

def _elements_to_markdown(elements) -> str:
    """Convert unstructured elements to Markdown."""
    from unstructured.documents.elements import (
        Title, ListItem, Table, FigureCaption,
        Image as ImageElement, Header, Footer, PageBreak,
    )

    md_parts: list[str] = []
    current_page = 0

    for el in elements:
        page = getattr(el.metadata, "page_number", None)
        if page and page != current_page:
            if current_page > 0:
                md_parts.append("")
            md_parts.append(f"--- Trang {page} ---")
            current_page = page

        if isinstance(el, Title):
            depth = getattr(el.metadata, "category_depth", None) or 1
            depth = min(max(int(depth), 1), 6)
            md_parts.append(f"\n{'#' * depth} {el.text}\n")
        elif isinstance(el, Table):
            html = getattr(el.metadata, "text_as_html", None)
            if html:
                md_parts.append(f"\n{html}\n")
            else:
                md_parts.append(f"\n```\n{el.text}\n```\n")
        elif isinstance(el, ListItem):
            md_parts.append(f"- {el.text}")
        elif isinstance(el, FigureCaption):
            md_parts.append(f"\n> **Hình:** {el.text}\n")
        elif isinstance(el, ImageElement):
            caption = (el.text or "").strip()
            if caption:
                md_parts.append(f"\n[Biểu đồ/Hình ảnh: {caption}]\n")
            else:
                md_parts.append("\n[Hình ảnh]\n")
        elif isinstance(el, (Header, Footer, PageBreak)):
            continue
        else:
            text = (el.text or "").strip()
            if text:
                md_parts.append(text)

    return "\n\n".join(md_parts)


def _elements_to_outline(elements) -> list[dict]:
    """Build outline from Title elements."""
    from unstructured.documents.elements import Title

    outline: list[dict] = []
    for el in elements:
        if isinstance(el, Title):
            depth = getattr(el.metadata, "category_depth", None) or 1
            page = getattr(el.metadata, "page_number", None) or 0
            outline.append({
                "level": min(max(int(depth), 1), 6),
                "title": el.text.strip()[:150],
                "page": page,
            })
    return outline


def _extract_with_unstructured(file_path: str) -> tuple[str, list[dict]]:
    """Smart extraction using the unstructured library."""
    from unstructured.partition.auto import partition

    elements = partition(
        filename=file_path,
        include_page_breaks=True,
        strategy="auto",
    )

    text = _elements_to_markdown(elements)
    outline = _elements_to_outline(elements)

    if not outline:
        outline = _auto_outline(text)

    return text.strip(), outline


# ---------------------------------------------------------------------------
# Custom PPTX extractor — extracts chart labels, series data, tables
# ---------------------------------------------------------------------------

def _extract_pptx(file_path: str) -> tuple[str, list[dict]]:
    """Extract slides to Markdown with chart label & table support."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    md_parts: list[str] = []
    outline: list[dict] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_texts: list[str] = []

        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text_frame.text.strip()

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                rows: list[str] = []
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        rows.append("| " + " | ".join(["---"] * len(row.cells)) + " |")
                if rows:
                    slide_texts.append("\n".join(rows))

            try:
                if hasattr(shape, "has_chart") and shape.has_chart:
                    chart = shape.chart
                    chart_parts: list[str] = []

                    if chart.has_title:
                        chart_parts.append(f"**{chart.chart_title.text_frame.text}**")

                    try:
                        cats = [str(c) for c in chart.plots[0].categories]
                        if cats:
                            chart_parts.append(f"Labels: {', '.join(cats)}")
                    except Exception:
                        pass

                    for series in chart.series:
                        name = series.name or "Series"
                        try:
                            vals = [str(v) for v in series.values]
                            chart_parts.append(f"{name}: {', '.join(vals)}")
                        except Exception:
                            pass

                    if chart_parts:
                        slide_texts.append("[Biểu đồ]\n" + "\n".join(chart_parts))
                    else:
                        slide_texts.append("[Biểu đồ]")
            except Exception:
                pass

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_texts.append("[Hình ảnh]")
            except Exception:
                pass

        if not slide_title:
            slide_title = f"Slide {i}"

        md_parts.append(f"## Slide {i}: {slide_title}\n\n" + "\n\n".join(slide_texts))
        outline.append({"level": 2, "title": f"Slide {i}: {slide_title}", "page": i})

    if not outline:
        outline.append({"level": 1, "title": "Toàn bộ tài liệu", "page": 0})

    return "\n\n".join(md_parts).strip(), outline


# ---------------------------------------------------------------------------
# Legacy fallback extractors
# ---------------------------------------------------------------------------

def _extract_pdf_legacy(file_path: str) -> tuple[str, list[dict]]:
    import fitz

    doc = fitz.open(file_path)
    full_text = ""
    outline: list[dict] = []

    for i, page in enumerate(doc):
        full_text += f"\n--- Trang {i + 1} ---\n{page.get_text()}"

    toc = doc.get_toc()
    for level, title, page in toc:
        outline.append({"level": level, "title": title, "page": page})

    doc.close()
    if not outline:
        outline = _auto_outline(full_text)
    return full_text.strip(), outline


def _extract_docx_legacy(file_path: str) -> tuple[str, list[dict]]:
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    full_text = ""
    outline: list[dict] = []

    for para in doc.paragraphs:
        full_text += para.text + "\n"
        style = para.style.name if para.style else ""
        if "Heading" in style:
            try:
                level = int(style.replace("Heading ", "").strip())
            except ValueError:
                level = 1
            outline.append({"level": level, "title": para.text.strip(), "page": 0})

    if not outline:
        outline = _auto_outline(full_text)
    return full_text.strip(), outline


def _extract_txt(file_path: str) -> tuple[str, list[dict]]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return text.strip(), _auto_outline(text)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _auto_outline(text: str) -> list[dict]:
    """Generate a basic outline when none exists."""
    lines = text.split("\n")
    outline: list[dict] = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        md_match = re.match(r"^(#{1,6})\s+(.+)", stripped)
        if md_match:
            outline.append({
                "level": len(md_match.group(1)),
                "title": md_match.group(2)[:150],
                "page": 0,
            })
            continue
        if re.match(r"^(Chương|Chapter|Phần|Bài|Mục|Section)\s", stripped, re.IGNORECASE):
            outline.append({"level": 1, "title": stripped[:100], "page": 0})
        elif re.match(r"^[IVXLCDM]+\.", stripped) or re.match(r"^\d+\.\s+[A-ZĐ]", stripped):
            outline.append({"level": 2, "title": stripped[:100], "page": 0})
    if not outline:
        outline.append({"level": 1, "title": "Toàn bộ tài liệu", "page": 0})
    return outline


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_text(file_path: str, file_type: str) -> tuple[str, list[dict]]:
    """Route to the best extractor for the given file type.

    Strategy:
      - PPTX/PPT  → custom extractor (chart labels + slide data)
      - PDF / DOCX → unstructured (smart layout detection, figure captions)
      - TXT        → direct read
    Falls back to legacy extractors when unstructured is unavailable.
    """
    ext = file_type.lower()

    if ext in ("pptx", "ppt"):
        try:
            return _extract_pptx(file_path)
        except ImportError:
            logger.warning("python-pptx not installed; trying unstructured")
        except Exception as e:
            logger.warning("Custom PPTX extraction failed: %s; trying unstructured", e)
        try:
            return _extract_with_unstructured(file_path)
        except Exception as e:
            logger.error("All PPTX extractors failed: %s", e)
            raise

    if ext in ("txt", "text/plain"):
        return _extract_txt(file_path)

    try:
        logger.info("Using unstructured parser for %s (%s)", file_path, ext)
        return _extract_with_unstructured(file_path)
    except ImportError:
        logger.warning("unstructured not installed; using legacy parser")
    except Exception as e:
        logger.warning("unstructured failed for %s: %s — falling back", file_path, e)

    if ext in ("pdf", "application/pdf"):
        return _extract_pdf_legacy(file_path)
    elif ext in ("docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
        return _extract_docx_legacy(file_path)
    else:
        return _extract_txt(file_path)


def chunk_text(text: str) -> list[str]:
    """Split text into chunks for embedding."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)
